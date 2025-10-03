from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# === INPUT DATA ===
song_title = "所有的荣耀归于祢"  # Replace with your song title
lyrics = """
荣耀归于祢 权柄归于祢
羔羊已作王 从今时到永远
荣耀归于祢 权柄归于祢
愿颂赞永远不停息

愿颂赞永远不停息
愿颂赞永远不停息

所有的荣耀归于祢
万民敬拜祢座前
所有的荣耀归于祢
荣耀荣耀主耶稣

所有的荣耀归于祢
万民敬拜祢座前
所有的荣耀归于祢
荣耀荣耀主耶稣

荣耀归于祢　权柄归于祢
羔羊已作王　从今时到永远
荣耀归于祢　权柄归于祢
愿颂赞永远不停息

荣耀归于祢 权柄归于祢
荣耀归于祢 权柄归于祢
荣耀归于祢 
权柄归于祢

所有的荣耀归于祢
万民敬拜祢座前
所有的荣耀归于祢
荣耀荣耀主耶稣

所有的荣耀归于祢
万民敬拜祢座前
所有的荣耀归于祢
荣耀荣耀主耶稣
荣耀荣耀主耶稣
荣耀荣耀主耶稣"""  # Replace with your entire lyrics (Chinese)

# Split lyrics into lines and group them into pairs
lines = [line.strip() for line in lyrics.split("\n") if line.strip()]
lyric_pairs = [lines[i:i+2] for i in range(0, len(lines), 2)]

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_font(run, font_name, font_size, bold=True, color=(0,0,0)):
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)

for idx, pair in enumerate(lyric_pairs, start=1):
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # --- Lyrics text box (slightly above vertical center) ---
    textbox_width = Inches(11.3)
    textbox_height = Inches(3)
    left = (prs.slide_width - textbox_width) / 2
    top = (prs.slide_height - textbox_height) / 2 - Inches(0.5)

    txBox = slide.shapes.add_textbox(left, top, textbox_width, textbox_height)
    tf = txBox.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Initial font size
    font_size = 60
    min_size = 30

    # Add paragraphs
    for line in pair:
        p = tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = Pt(font_size * 2 / 1.2)
        run = p.runs[0]
        set_font(run, "KaiTi", font_size, True, (0,0,0))

    # Shrink font until text fits the box
    while txBox.height > textbox_height and font_size > min_size:
        font_size -= 1
        for p in tf.paragraphs:
            for run in p.runs:
                set_font(run, "KaiTi", font_size, True, (0,0,0))

    # --- Bottom left: Song Title ---
    txBox_title = slide.shapes.add_textbox(Inches(0.3), Inches(6.7), Inches(5), Inches(0.8))
    tf_title = txBox_title.text_frame
    tf_title.text = song_title
    tf_title.paragraphs[0].alignment = PP_ALIGN.LEFT
    run_title = tf_title.paragraphs[0].runs[0]
    set_font(run_title, "KaiTi", 18, True, (255, 192, 0))

    # --- Bottom center: Page Count ---
    txBox_page = slide.shapes.add_textbox(Inches(6), Inches(6.7), Inches(1.5), Inches(0.8))
    tf_page = txBox_page.text_frame
    tf_page.text = f"{idx}/{len(lyric_pairs)}"
    tf_page.paragraphs[0].alignment = PP_ALIGN.CENTER
    run_page = tf_page.paragraphs[0].runs[0]
    set_font(run_page, "KaiTi", 18, True, (255,255,255))

# Save presentation
prs.save("song_lyrics.pptx")
