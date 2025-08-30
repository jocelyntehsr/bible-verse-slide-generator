import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


# === Load Excel ===
# Columns: Title (CN), Title (EN), Number, Chinese Verse, English Verse
df = pd.read_csv("bible_verses.csv")
eng_df = pd.read_csv("bible_verses_en.csv")

# === Create Presentation ===
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === Font helper ===
def set_font(run, font_name, font_size, bold=False, color=RGBColor(0,0,0)):
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color

# === Add textbox with shrinking font ===
def add_textbox_auto_font(slide, left, top, width, height, text, font_name, max_font_size, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
#    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)

    font_size = max_font_size
    set_font(run, font_name, font_size, bold)

    # Reduce font size until text fits vertically
    while tb.height > Inches(height) and font_size > 10:
        font_size -= 1
        set_font(run, font_name, font_size, bold)

    return tb

# === Create slides ===
for (_, row1), (_,row2) in zip(df.iterrows(), eng_df.iterrows()):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    chinese_title = f"{row1['Book Name']} {row1['Chapter']}:{row1['Verse']}"
    english_title = f"{row2['Book Name']} {row2['Chapter']}:{row2['Verse']}"

    # Background black
#    background_image_path = "background.jpg"
#    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
#    fill = slide.background.fill
#    fill.user_picture(background_image_path)

    # 1. Title bar split into 2 boxes (CN left, EN right)
    add_textbox_auto_font(slide, 0.25, 0.2, 4, 0.5, chinese_title, "KaiTi", 20, bold=True, align=PP_ALIGN.LEFT)
    add_textbox_auto_font(slide, 0.25, 0.7, 4, 0.5, english_title, "Calibri", 20, bold=True, align=PP_ALIGN.LEFT)

    # 2. Chinese verse
    add_textbox_auto_font(slide, 0.25, 1.45, 12.8, 2.8, row1['CleanText'], "KaiTi", 50, bold=True, align=PP_ALIGN.LEFT)

    # 3. English verse
    add_textbox_auto_font(slide, 0.25, 4.25, 12.8, 2.75, row2['CleanText'], "Calibri", 35, bold=True, align=PP_ALIGN.LEFT)

# === Save ===
prs.save("WW35_main_slide.pptx")
print("✅ Presentation created: Main_Slide.pptx")
