import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

# === Load Excel ===
df1 = pd.read_csv("bible_verses.csv")
df2 = pd.read_csv("bible_verses_en.csv")

# === Create Presentation ===
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# === Font helper ===
def set_font(run, font_name, font_size, bold=False, color=RGBColor(255,255,255)):
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color

# === Add textbox with shrinking font ===
def add_textbox_auto_font(slide, left, top, width, height, text, font_name, max_font_size, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
#    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)

    font_size = max_font_size
    set_font(run, font_name, font_size, bold)

    # Reduce font size until text fits vertically
    while tb.height > Inches(height) and font_size > 10:
        font_size -= 1
        set_font(run, font_name, font_size, bold)

    return tb

# === Create slides ===
for (_, row1), (_,row2) in zip(df1.iterrows(), df2.iterrows()):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background black
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    # 1. Title bar split into 2 boxes (CN left, EN right)
    add_textbox_auto_font(slide, 0, 0, 3.5, 0.5, row1['Title (CN)'], "KaiTi", 20, bold=True, align=PP_ALIGN.LEFT)
    add_textbox_auto_font(slide, 3.5, 0, 4.0, 0.5, row1['Title (EN)'], "Calibri", 20, bold=True, align=PP_ALIGN.LEFT)

    # 2. Number
    add_textbox_auto_font(slide, 0, 0.5, 0.5, 1.625, row1['Verse'], "Calibri", 20, bold=True, align=PP_ALIGN.CENTER)
    add_textbox_auto_font(slide, 0, 2.175, 0.5, 1.625, row1['Verse'], "Calibri", 20, bold=True, align=PP_ALIGN.CENTER)

    # 3. Chinese verse
    add_textbox_auto_font(slide, 0.5, 0.5, 9.5, 1.625, row1['CleanText'], "KaiTi", 30, bold=True, align=PP_ALIGN.LEFT)

    # 4. English verse
    add_textbox_auto_font(slide, 0.5, 2.125, 9.5, 1.625, row2['CleanText'], "Calibri", 25, bold=True, align=PP_ALIGN.LEFT)

# === Save ===
prs.save("ww39_圣经经节.pptx")
print("✅ Presentation created: 圣经经节.pptx")
